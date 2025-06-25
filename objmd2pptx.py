from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup, NavigableString
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.simpletypes import ST_Percentage
import os
from pathlib import Path
import md2docx as m2d
import markdown2    
from urllib.parse import urlparse


class MD2PPTX:
    def __init__(self, path_or_content, out_pptx_file_path:Path):
        self.content = path_or_content
        self.out_pptx_file_path = out_pptx_file_path
        self.prs = Presentation()
        self.md_path = Path()
        if isinstance(path_or_content, Path):
            self.md_path = path_or_content
            self.content = path_or_content.read_text(encoding='utf-8')
        elif isinstance(path_or_content, str):
            self.content = path_or_content
        
    
    def convert(self):
        self.md_text = '\n'.join(line for line in self.content.splitlines() if line.strip())  # 删除空行
        self.html = markdown2.markdown(self.md_text, extras=[ "fenced-code","toc","tables"])  #
    
        self.soup = BeautifulSoup(self.html, 'lxml')
        self.body = self.soup.body if self.soup.body else self.soup

        for child in self.body.children:
            if hasattr(child, 'name'):
                #process_element(md_path,child,doc)
                print(child)    
            elif isinstance(child, NavigableString):
                text = child.strip()
                if text:
                    print(text)
        self.save_pptx()

    def process_element(self, element, doc, main_p=None):
        # 处理块级标签（标题、段落、div、列表、表格等）。
        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            level = int(element.name[1])
            p = doc.add_heading("", level=level)
        
        if element.get('align') == 'center':
            p.alignment = PP_ALIGN.CENTER




    def add_slide(self, title):
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

    def add_text(self, text):
        text_frame = self.prs.slides[-1].shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.0), Inches(0.8))
        text_frame.text = text

    def save_pptx(self):
        self.prs.save(self.pptx_file_path)

if __name__ == '__main__':
    md_file_path = 'f:\DevProject\pt\md2\example.md'
    pptx_file_path = 'f:\DevProject\pt\md2\example.pptx'
    md2pptx = MD2PPTX(md_file_path, pptx_file_path)
    md2pptx.convert()    