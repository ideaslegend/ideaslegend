import markdown2
from bs4 import BeautifulSoup, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import md2docx as M2D
import threading



# 设置每张幻灯片文本框所可容纳的最大字符数（根据需要调整）
A4_WIDTH = Cm(21)  
A4_HEIGHT = Cm(29.7)  
MAX_CHARS_PER_SLIDE = 800  # 根据 A4 纵向页面，适当调整最大字符数

def split_text(text, max_chars):
    """
    将文本按空格拆分为多个段落，保证每个段落的长度不超过 max_chars。
    尽可能保持单词完整。
    """
    words = text.split()
    chunks = []
    current_chunk = ""
    for word in words:
        # 加一个空格后 word 长度
        if len(current_chunk) + len(word) + 1 > max_chars:
            chunks.append(current_chunk.strip())
            current_chunk = word + " "
        else:
            current_chunk += word + " "
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

def adjust_font_size(text, max_chars):
    """根据文本长度动态调整字体大小"""
    
    base_font_size = 36
    min_font_size = 8
    if len(text) == 0:
        return min_font_size
    scaling_factor = max_chars / len(text)
    adjusted_size = int(base_font_size * scaling_factor)
    return max(min_font_size, min(adjusted_size, base_font_size))

def process_table_element(table_element, prs):
    """
    解析 HTML 表格，将其转换为 PPTX 表格。
    每个 <table> 元素单独生成一页幻灯片，自动解析表格的行和列，
    并为每个单元格设置文本及自动换行。
    """
    # 解析每一行的数据；支持<th>和<td>
    rows_data = []
    for tr in table_element.find_all("tr"):
        row = []
        for cell in tr.find_all(['th', 'td']):
            row.append(cell)
        if row:
            rows_data.append(row)
    
    if not rows_data:
        return

    # 取得表格行数和最大列数：
    num_rows = len(rows_data)
    num_cols = max(len(row) for row in rows_data)

    # 创建新的幻灯片，并添加表格形状
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 设定边距（可以适当调整）
    margin_x = Inches(0.5)  # 左右边距

    # 计算表格尺寸
    width = prs.slide_width - 2 * margin_x  # 让表格填充可用宽度
    height = Inches(3)  # 让表格填充可用高度

    # 计算居中位置
    left = Inches(0.5)
    top = Inches(0.5)

    # 创建表格并居中
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    
    for i, row in enumerate(rows_data):
        for j in range(num_cols):
            cell = table.cell(i, j)
            try:
                cell.text_frame.word_wrap = True
            except Exception:
                pass  # 如果版本不支持设置 word_wrap，可忽略

            if j < len(row):
                cell_element = row[j]
                # 使用 cell.text_frame 的第一个段落（先清空已有内容）
                tf = cell.text_frame
                if tf.paragraphs:
                    if 'style' in cell_element.attrs:
                        align_map = {
                            'text-align:left;': PP_ALIGN.LEFT,
                            'text-align:center;': PP_ALIGN.CENTER,
                            'text-align:right;': PP_ALIGN.RIGHT,
                            'text-align:justify;' : PP_ALIGN.JUSTIFY
                        }
                        tf.paragraphs[0].alignment = align_map.get(cell_element['style'].lower(), PP_ALIGN.LEFT)
                    p = tf.paragraphs[0]
                    p.text = ""  # 清空原始文本
                else:
                    p = tf.add_paragraph()
                process_inline_element(prs,cell_element, p)
            else:
                cell.text = ""
def process_block_element(presentation,element, text_frame, bullet_level=0):
    """
    根据块级标签生成文本框中的段落，支持 p, h1-h6, blockquote, ul/ol 等
    bullet_level 用于处理列表嵌套层级
    """
    tag = element.name.lower()
    
     # 对 p 与 div 元素先判断是否需要分页（注意 div 也可能包含大量文本）
    if tag in ['p', 'div']:
        #plain_text = element.get_text(separator=" ", strip=True)
        plain_text = element.get_text()
        if len(plain_text) > MAX_CHARS_PER_SLIDE:
            # 如果文字太多，则按 max_chars 拆分，并为每个块新建幻灯片
            prs = presentation
            margin_x = Inches(0.5)  # 左右边距
            width = prs.slide_width - 2 * margin_x
            chunks = split_text(plain_text, MAX_CHARS_PER_SLIDE)
            for chunk in chunks:
                new_slide = prs.slides.add_slide(prs.slide_layouts[6])
                tx_box = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), width, Inches(5))
                tx_box.text_frame.word_wrap = True
                para = tx_box.text_frame.add_paragraph()
                para.text = chunk
                para.font.size = Pt(adjust_font_size(para.text, MAX_CHARS_PER_SLIDE))
            return  # 不再继续对该 element 进行内联处理
    # 如果是标题或段落
    if tag in ['div','p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'blockquote']:
        p = text_frame.add_paragraph()
        p.level = bullet_level
        # 针对标题可以调整字号
        if tag.startswith('h'):
            try:
                level = int(tag[1])
            except ValueError:
                level = 1
            # 举例：h1 使用 40pt，后续依次减小字号
            p.font.size = Pt(max(40 - (level - 1) * 2, 20))
        # 对 blockquote 可设置为斜体或者其他格式（这里作为示例）
        if tag == 'blockquote':
            p.font.italic = True

        process_inline_element(presentation,element, p)

    # 如果是无序列表或有序列表
    elif tag in ['ul', 'ol']:
        for li in element.find_all('li', recursive=False):
            p = text_frame.add_paragraph()
            p.level = bullet_level
            p.font.size = Pt(adjust_font_size(element.get_text(), MAX_CHARS_PER_SLIDE))
            process_inline_element(presentation,li, p)
            # 检查 li 内是否嵌套了子列表
            for sublist in li.find_all(['ul', 'ol'], recursive=False):
                process_block_element(presentation,sublist, text_frame, bullet_level=bullet_level+1)
    elif tag == 'table':
                process_table_element(element, presentation)
    elif tag == 'img':
                # 顶级图片元素处理
                process_img_element(element, presentation)
    else:
        # 其他情况，统一当作一段处理
        p = text_frame.add_paragraph()
        p.font.size = Pt(adjust_font_size(element.get_text(), MAX_CHARS_PER_SLIDE))
        process_inline_element(presentation,element, p)

def process_inline_element(presentation,element, paragraph):
    """
    解析 element 内的内容，将内联元素（包括超链接等）添加到段落中。
    如果遇到嵌套标记则递归处理，支持 a, strong/b, em/i, code 等
    """
    for child in element.contents:
        if isinstance(child, NavigableString):
            text = child.string
            if text and text.strip():
                run = paragraph.add_run()
                run.text = text
                #run.font.size = Pt(adjust_font_size(run.text, MAX_CHARS_PER_SLIDE))
        elif child.name:
            tag = child.name.lower()
            if tag == 'a':
                run = paragraph.add_run()
                run.text = child.get_text()
                href = child.get('href')
                #run.font.name = 'Courier New'
                #run.font.size = Pt(adjust_font_size(run.text, MAX_CHARS_PER_SLIDE))
                if href:
                    # 设定超链接（注意：python-pptx 对此支持有限，可能在不同版本中表现不同）
                    run.hyperlink.address = href
            elif tag in ['strong', 'b']:
                run = paragraph.add_run()
                run.text = child.get_text()
                run.font.bold = True
                #run.font.size = Pt(adjust_font_size(run.text, MAX_CHARS_PER_SLIDE))
                process_inline_element(presentation,child, paragraph)
            elif tag in ['em', 'i']:
                run = paragraph.add_run()
                run.text = child.get_text()
                run.font.italic = True
                #run.font.size = Pt(adjust_font_size(run.text, MAX_CHARS_PER_SLIDE))
                process_inline_element(presentation,child, paragraph)
            elif tag == 'code':
                
                # # 创建带有灰色填充的文本框（模拟代码块背景）
                # tx_box = paragraph.shapes.add_textbox(Cm(2), Cm(3), Cm(17), Cm(1))  # 设置合适的位置和大小
                # tx_box.fill.solid()
                # tx_box.fill.fore_color.rgb = RGBColor(200, 200, 200)  # 设置灰色背景

                # tf = tx_box.text_frame
                # tf.word_wrap = True
                # p = tf.add_paragraph()
                
                # run = p.add_run()
                # run.text = child.get_text()
                # run.font.name = 'Courier New'

                run = paragraph.add_run()
                run.text = child.get_text()
                #run.font.size = Pt(adjust_font_size(run.text, MAX_CHARS_PER_SLIDE))
                # 可选：设置等宽字体，便于展示代码
                run.font.name = 'Courier New'
            elif tag == 'img':
                # 内联图片：先在当前段落插入占位符
                alt = child.get('alt', '')
                src = child.get('src', '')
                placeholder = f"[图片: {alt}]" if alt else "[图片]"
                run = paragraph.add_run()
                run.text = placeholder
                process_img_element(child, presentation)
            else:
                # 对于其他标签，比如 span、u 等直接递归处理
                process_inline_element(presentation,child, paragraph)

def process_img_element(img_element, prs):
    """
    解析 HTML 图片标签，将其转换为 PPTX 图片。
    每个图片元素单独生成一页幻灯片，通过图片的 src 属性添加图片，
    同时如果存在 alt 属性，则在图片下方添加说明文字。
    """
    src = img_element.get('src', '')

    if src.startswith('http://') or src.startswith('https://'):
        import queue
        def thread_function(queue):
            path = M2D.download_image(src,M2D.get_current_file_dir()+ '/temp'+'/images')
            queue.put(path)  
        q = queue.Queue()
        thread = threading.Thread(target=thread_function,args=(q,))
        thread.start()
        thread.join(5)
        img_path = q.get() if not q.empty() else None
        
        if not img_path:
            print(f"下载失败{src}")
                
    else:
        md_path = Path(M2D.get_current_file_dir()).resolve()
        img_path = (md_path/src).resolve()
    if not src:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin_x = Inches(0.5)  # 左右边距
    
    width = prs.slide_width - 2 * margin_x  
    left = Inches(0.5)
    top = Inches(0.5)
    
    try:
        if str(img_path).endswith('.svg'):
                    png_path = img_path.rsplit('.', 1)[0] + '.png'
                    M2D.svg_to_png(img_path, png_path)
                    img_path = png_path
        pic = slide.shapes.add_picture(str(img_path), left, top, width=width)
        alt = img_element.get('alt', '')
        if alt:
            # 在图片下方添加说明文字（caption）
            caption_top = top + pic.height + Inches(0.2)
            tx_box = slide.shapes.add_textbox(left, caption_top, width, Inches(1))
            tx_box.text_frame.word_wrap = True
            tx_box.text_frame.text = alt
    except Exception as e:
        print(f"无法添加图片 {src}: {e}")




def md_to_pptx(md_file_path:Path, pptx_file_path:Path):
    """
    读取 Markdown 文件，转换成 HTML，再解析生成 PPTX 文件
    """
    md_content =md_file_path.read_text(encoding='utf-8')
    md_content = '\n'.join(line for line in md_content.splitlines() if line.strip())  # 删除空行
    html = markdown2.markdown(md_content,extras=["fenced-code","toc","tables"])
    soup = BeautifulSoup(html, 'lxml')
    prs = Presentation()
    prs.slide_width = A4_WIDTH
    prs.slide_height = A4_HEIGHT
    
    #elements = soup.body.contents if soup.body else soup.contents
    elements = soup.body if soup.body else soup
    margin_x = Inches(0.5)  # 左右边距
    width = prs.slide_width - 2 * margin_x

    #for elem in elements:
    for elem in elements.children:
        # 忽略空白文本
        if isinstance(elem, NavigableString):
            if not elem.strip():
                continue
            # 如果有裸文本，则也生成一页幻灯片来显示
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), width, Inches(5))
            tf = tx_box.text_frame
            p = tf.add_paragraph()
            p.text = elem.strip()
            p.font.size = Pt(adjust_font_size(p.text, MAX_CHARS_PER_SLIDE))
        elif elem.name:
            # 为每个块级元素创建一页幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), width, Inches(5))
            tf = tx_box.text_frame
            process_block_element(prs,elem, tf)
            tx_box.text_frame.word_wrap = True  # 设置自动换行

    # 保存 PPTX 到指定文件
    prs.save(pptx_file_path)
    print(f"PPTX 文件已保存：{pptx_file_path}")
if __name__ == '__main__':
    # 示例：假设 md 文件名为 "example.md"，生成的 PPT 文件名为 "output.pptx"
    md_file = Path("./common/test1.md")
    pptx_file = Path("./common/readme_2.pptx")
    md_to_pptx(md_file, pptx_file)
